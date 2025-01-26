import os
import io
import re
import logging
import docx
import pptx
import pandas as pd
import pytesseract
from PIL import Image
from io import BytesIO
from PyPDF2 import PdfReader

import google.generativeai as genai

import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords

from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup
)
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    CallbackContext,
    ContextTypes,
    filters
)
from geopy.geocoders import Nominatim
from dotenv import load_dotenv

# NLTK data
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('averaged_perceptron_tagger')
nltk.download('maxent_ne_chunker')
nltk.download('words')

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO
)
logger = logging.getLogger(__name__)

load_dotenv()

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
CHEMIST_CHAT_ID = os.getenv("CHEMIST_CHAT_ID", "")

###############################################################################
# Configure Gemini
###############################################################################
genai.configure(api_key=GOOGLE_API_KEY)
vision_model = genai.GenerativeModel("gemini-pro-vision")
llm_model = genai.GenerativeModel("gemini-pro")

###############################################################################
# OCR / Text Extraction
###############################################################################
def extract_text_from_file(file_bytes: bytes, file_extension: str) -> str:
    """
    1) For PDFs, docx, pptx, CSV -> extract text directly
    2) For images -> first try gemini-pro-vision. If it fails, fallback to local Tesseract.
    """
    text = ""

    # PDF
    if file_extension == ".pdf":
        pdf_io = BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_io)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text

    # Images
    elif file_extension in [".png", ".jpg", ".jpeg", ".webp"]:
        img_file = BytesIO(file_bytes)
        img = Image.open(img_file)

        # Try gemini-pro-vision first
        try:
            response = vision_model.generate_content(
                [f"Extract textual information from this picture:", img]
            )
            extracted = response.text.strip()
            # If Gemini returns empty or triggers a refusal, fallback
            if extracted:
                text += extracted
            else:
                raise ValueError("Gemini content empty or refused.")
        except Exception as e:
            logger.error(f"Gemini Vision fallback triggered: {e}")
            # Fallback to local OCR
            text += pytesseract.image_to_string(img)

    # docx
    elif file_extension == ".docx":
        doc = docx.Document(BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"

    # pptx
    elif file_extension == ".pptx":
        presentation = pptx.Presentation(BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    # csv
    elif file_extension == ".csv":
        df = pd.read_csv(BytesIO(file_bytes))
        text += df.to_string()

    return text.strip()

###############################################################################
# Extract Medicines with Gemini
###############################################################################
def extract_medicines_gemini(prescription_text: str) -> list:
    """
    Uses gemini-pro to parse the text and extract medicine names.
    If gemini refuses or returns nothing, return an empty list.
    """
    if not prescription_text:
        return []

    prompt = f"""
Extract all medicine names from the text below, with only one medicine per line, and no extra explanations.

Prescription:
{prescription_text}
"""
    try:
        response = llm_model.generate_content(prompt)
        lines = response.text.strip().split("\n")
        meds = [m.strip() for m in lines if m.strip()]
        return meds
    except Exception as e:
        logger.error(f"Gemini extraction error: {e}")
        return []

###############################################################################
# Billing
###############################################################################
def generate_bill(meds: list, days: int, address: str) -> str:
    """
    Format a creative, neat prescription bill summary.
    Each medicine has {name, freq}.
    """
    total_qty = 0
    result = "🧾 **Awesome Prescription Summary**\n\n"
    for med in meds:
        qty = med["freq"] * days
        total_qty += qty
        result += (
            f"💊 **{med['name']}**\n"
            f"   • Daily Frequency: {med['freq']}\n"
            f"   • Duration (days): {days}\n"
            f"   • Total Quantity: {qty}\n\n"
        )
    result += f"**Total Medicines:** {len(meds)}\n"
    result += f"**Combined Quantity:** {total_qty}\n"
    if address:
        result += f"📍 **Delivery Address**: {address}\n"
    return result

###############################################################################
# Bot Class
###############################################################################
class RxConnectTeleBot:
    """
    - /start: greeting
    - Upload files to parse prescription
    - Extract medicine names automatically
    - Ask user if they want to set freq or skip
    - If set freq -> inline keyboard for each medicine
    - /setdays <X> -> user sets how many days
    - /setaddress <Address> or share location
    - /getbill -> finalize
    """
    def __init__(self):
        # user_data: { user_id: { "prescription_text": "", "meds": [{"name":"...", "freq":0}, ...], "days":0, "address":""} }
        self.user_data = {}

    def get_state(self, user_id: int):
        if user_id not in self.user_data:
            self.user_data[user_id] = {
                "prescription_text": "",
                "meds": [],
                "days": 0,
                "address": ""
            }
        return self.user_data[user_id]

    async def start(self, update: Update, context: CallbackContext) -> None:
        await update.message.reply_text(
            "👋 Hello! I'm your creative AI Prescription Assistant. "
            "1) Send me a PDF/Image/Doc of your prescription.\n"
            "2) I'll extract medicine names.\n"
            "3) You can set how many times a day you take each.\n"
            "4) Then set how many days.\n"
            "5) Get a neat final bill!\n\n"
            "Ready when you are!"
        )

    async def handle_prescription(self, update: Update, context: CallbackContext) -> None:
        user_id = update.effective_user.id
        state = self.get_state(user_id)

        file_bytes = None
        file_ext = ".pdf"

        if update.message.document:
            file_ext = os.path.splitext(update.message.document.file_name)[1].lower()
            doc_file = await update.message.document.get_file()
            file_bytes = await doc_file.download_as_bytearray()
        elif update.message.photo:
            file_ext = ".jpg"
            photo = update.message.photo[-1]
            doc_file = await photo.get_file()
            file_bytes = await doc_file.download_as_bytearray()

        if not file_bytes:
            await update.message.reply_text("I didn't receive any file. Please try again.")
            return

        await update.message.reply_text("⏳ Processing your file...")
        text = extract_text_from_file(file_bytes, file_ext)
        state["prescription_text"] = text

        meds_raw = extract_medicines_gemini(text)
        state["meds"] = [{"name": m, "freq": 0} for m in meds_raw]

        if not meds_raw:
            await update.message.reply_text(
                "Hmm, I didn't find any medicines. "
                "You can still try again or manually share details."
            )
            return

        # Ask user if they want to set frequencies
        keyboard = [
            [InlineKeyboardButton("Yes, let me set frequencies", callback_data="freq_yes")],
            [InlineKeyboardButton("No, skip frequencies", callback_data="freq_no")]
        ]
        await update.message.reply_text(
            "Detected these medicines:\n" + "\n".join(f"• {m}" for m in meds_raw) +
            "\nDo you want to set daily frequency for each medicine?",
            reply_markup=InlineKeyboardMarkup(keyboard)
        )

    async def process_freq_choice(self, update: Update, context: CallbackContext):
        query = update.callback_query
        await query.answer()

        user_id = query.from_user.id
        state = self.get_state(user_id)
        choice = query.data  # freq_yes or freq_no

        if choice == "freq_yes":
            await query.edit_message_text("Great! Let's set frequencies one by one.")
            # Start asking frequencies for each uncategorized med
            await self.ask_freq_for_next(query.message, user_id, context)
        else:
            # skip freq -> default freq=1 for all
            for med in state["meds"]:
                med["freq"] = 1
            await query.edit_message_text("Skipping frequency. Each medicine set to 1 time a day by default.")
            await query.message.reply_text("Now set days using /setdays <X>.")

    async def ask_freq_for_next(self, msg, user_id: int, context: CallbackContext):
        state = self.get_state(user_id)
        med_list = state["meds"]

        for i, med in enumerate(med_list):
            if med["freq"] == 0:
                buttons = [
                    [
                        InlineKeyboardButton("1 Time", callback_data=f"set_{i}_1"),
                        InlineKeyboardButton("2 Times", callback_data=f"set_{i}_2"),
                        InlineKeyboardButton("3 Times", callback_data=f"set_{i}_3")
                    ]
                ]
                await msg.reply_text(
                    f"How many times a day for: {med['name']}?",
                    reply_markup=InlineKeyboardMarkup(buttons)
                )
                return

        # If no med with freq=0, ask for days
        await msg.reply_text("All frequencies set! Use /setdays <X> to specify how many days.")

    async def freq_setter(self, update: Update, context: CallbackContext):
        query = update.callback_query
        await query.answer()
        data = query.data  # e.g. set_0_2
        user_id = query.from_user.id
        state = self.get_state(user_id)

        parts = data.split("_")
        if len(parts) == 3:
            index = int(parts[1])
            freq_val = int(parts[2])
            if 0 <= index < len(state["meds"]):
                state["meds"][index]["freq"] = freq_val
                med_name = state["meds"][index]["name"]
                await query.edit_message_text(f"Set {med_name} to {freq_val} time(s) per day.")
                # Next frequency
                dummy_msg = query.message
                await self.ask_freq_for_next(dummy_msg, user_id, context)

    async def setdays(self, update: Update, context: CallbackContext):
        user_id = update.effective_user.id
        state = self.get_state(user_id)

        try:
            days_val = int(context.args[0])
            state["days"] = days_val
            await update.message.reply_text(
                f"👍 Duration set to {days_val} day(s).\n\nShare your address with /setaddress <address> "
                "or send your location so I can finalize."
            )
        except:
            await update.message.reply_text("Use /setdays <integer>.")

    async def setaddress(self, update: Update, context: CallbackContext):
        user_id = update.effective_user.id
        state = self.get_state(user_id)
        addr_args = context.args

        if not addr_args:
            await update.message.reply_text("Usage: /setaddress <your full address>")
            return

        address_str = " ".join(addr_args)
        state["address"] = address_str
        await update.message.reply_text("Address saved. Use /getbill to see your final prescription summary.")

    async def handle_location(self, update: Update, context: CallbackContext):
        user_id = update.effective_user.id
        state = self.get_state(user_id)

        loc = update.message.location
        if loc:
            geolocator = Nominatim(user_agent="RxConnectBot")
            location_info = geolocator.reverse(f"{loc.latitude}, {loc.longitude}", exactly_one=True)
            address_text = location_info.address if location_info else "Unknown"
            state["address"] = address_text
            await update.message.reply_text(f"Location saved as: {address_text}\nUse /getbill to finalize.")
        else:
            await update.message.reply_text("No valid location found.")

    async def getbill(self, update: Update, context: CallbackContext):
        user_id = update.effective_user.id
        state = self.get_state(user_id)

        meds = state["meds"]
        days_val = state["days"]
        address_val = state["address"]

        if not meds:
            await update.message.reply_text("No medicines found. Please upload a prescription first.")
            return
        if days_val <= 0:
            await update.message.reply_text("Days not set. Use /setdays <X> first.")
            return

        bill_text = generate_bill(meds, days_val, address_val)
        await update.message.reply_text(bill_text, parse_mode="Markdown")

        # Forward to chemist
        if CHEMIST_CHAT_ID:
            await context.bot.send_message(
                chat_id=CHEMIST_CHAT_ID,
                text=f"📢 New Prescription from User {user_id}:\n\n{bill_text}",
                parse_mode="Markdown"
            )

###############################################################################
# Main
###############################################################################
def main():
    bot = RxConnectTeleBot()
    application = ApplicationBuilder().token(TELEGRAM_BOT_TOKEN).build()

    # Commands
    application.add_handler(CommandHandler("start", bot.start))
    application.add_handler(CommandHandler("setdays", bot.setdays))
    application.add_handler(CommandHandler("setaddress", bot.setaddress))
    application.add_handler(CommandHandler("getbill", bot.getbill))

    # File processing
    application.add_handler(MessageHandler(filters.Document.ALL, bot.handle_prescription))
    application.add_handler(MessageHandler(filters.PHOTO, bot.handle_prescription))

    # Location
    application.add_handler(MessageHandler(filters.LOCATION, bot.handle_location))

    # Frequency decisions
    application.add_handler(CallbackQueryHandler(bot.process_freq_choice, pattern="^(freq_yes|freq_no)$"))
    application.add_handler(CallbackQueryHandler(bot.freq_setter, pattern="^set_"))

    application.run_polling()

if __name__ == "__main__":
    main()
